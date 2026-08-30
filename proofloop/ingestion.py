"""Safe, bounded extraction for Business Context evidence."""

from __future__ import annotations

import asyncio
import csv
import io
import ipaddress
import socket
from pathlib import Path
from typing import Any
from urllib.parse import urlparse

import httpx
from bs4 import BeautifulSoup


MAX_FILE_BYTES = 10 * 1024 * 1024
MAX_EXTRACTED_CHARS = 45_000
MAX_URL_BYTES = 1_500_000
SUPPORTED_EXTENSIONS = {".pdf", ".csv", ".xlsx", ".txt", ".md", ".docx", ".pptx"}


class EvidenceExtractionError(ValueError):
    """Raised when a source cannot be safely or usefully extracted."""


def _trim(text: str) -> str:
    compact = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    return compact[:MAX_EXTRACTED_CHARS]


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n".join((page.extract_text() or "") for page in reader.pages[:100])


def _extract_csv(content: bytes) -> str:
    decoded = content.decode("utf-8-sig", errors="replace")
    rows = csv.reader(io.StringIO(decoded))
    rendered: list[str] = []
    for index, row in enumerate(rows):
        if index >= 500:
            break
        rendered.append(" | ".join(cell.strip() for cell in row[:50]))
    return "\n".join(rendered)


def _extract_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    rendered: list[str] = []
    for sheet in workbook.worksheets[:5]:
        rendered.append(f"Sheet: {sheet.title}")
        for index, row in enumerate(sheet.iter_rows(values_only=True)):
            if index >= 250:
                break
            rendered.append(" | ".join("" if value is None else str(value) for value in row[:50]))
    return "\n".join(rendered)


def _extract_docx(content: bytes) -> str:
    from docx import Document

    document = Document(io.BytesIO(content))
    rendered = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables[:25]:
        for row in table.rows[:100]:
            rendered.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(rendered)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    rendered: list[str] = []
    for index, slide in enumerate(presentation.slides[:100], start=1):
        rendered.append(f"Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                rendered.append(shape.text)
    return "\n".join(rendered)


def extract_uploaded_file(filename: str, content: bytes) -> dict[str, Any]:
    """Extract text from one supported upload without retaining the source bytes."""

    if not filename:
        raise EvidenceExtractionError("Every upload must have a filename.")
    if len(content) > MAX_FILE_BYTES:
        raise EvidenceExtractionError(f"{filename} exceeds the 10 MB upload limit.")

    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise EvidenceExtractionError(
            f"{filename} is not supported. Use PDF, CSV, XLSX, TXT, Markdown, DOCX, or PPTX."
        )

    try:
        if extension == ".pdf":
            text = _extract_pdf(content)
        elif extension == ".csv":
            text = _extract_csv(content)
        elif extension == ".xlsx":
            text = _extract_xlsx(content)
        elif extension == ".docx":
            text = _extract_docx(content)
        elif extension == ".pptx":
            text = _extract_pptx(content)
        else:
            text = content.decode("utf-8-sig", errors="replace")
    except Exception as exc:  # parsers expose several format-specific errors
        raise EvidenceExtractionError(f"ProofLoop could not read {filename}.") from exc

    extracted = _trim(text)
    if not extracted:
        raise EvidenceExtractionError(f"No readable text was found in {filename}.")
    return {
        "name": filename,
        "extension": extension,
        "byte_count": len(content),
        "text": extracted,
        "source_note": "Original retained only in the user's browser; server bytes discarded after extraction.",
    }


def _validate_public_host(hostname: str) -> None:
    addresses = socket.getaddrinfo(hostname, 443, type=socket.SOCK_STREAM)
    for entry in addresses:
        address = ipaddress.ip_address(entry[4][0])
        if any(
            [
                address.is_private,
                address.is_loopback,
                address.is_link_local,
                address.is_multicast,
                address.is_reserved,
                address.is_unspecified,
            ]
        ):
            raise EvidenceExtractionError("Only public website URLs can be imported.")


async def extract_public_url(url: str) -> dict[str, Any]:
    """Fetch one public HTTPS page with SSRF and response-size protections."""

    parsed = urlparse(url.strip())
    if parsed.scheme != "https" or not parsed.hostname:
        raise EvidenceExtractionError("Website evidence must use a public HTTPS URL.")
    await asyncio.to_thread(_validate_public_host, parsed.hostname)

    async with httpx.AsyncClient(
        follow_redirects=False,
        timeout=httpx.Timeout(15.0),
        headers={"User-Agent": "ProofLoopEvidenceReader/1.0"},
    ) as client:
        response = await client.get(url)
    if response.status_code >= 400:
        raise EvidenceExtractionError(f"Website returned HTTP {response.status_code}.")
    if 300 <= response.status_code < 400:
        raise EvidenceExtractionError("Redirecting website URLs must be provided at their final public URL.")
    content_type = response.headers.get("content-type", "").lower()
    if "text/html" not in content_type and "text/plain" not in content_type:
        raise EvidenceExtractionError("The URL must return an HTML or plain-text page.")
    if len(response.content) > MAX_URL_BYTES:
        raise EvidenceExtractionError("The website page is too large to import safely.")

    if "text/html" in content_type:
        soup = BeautifulSoup(response.text, "html.parser")
        for element in soup(["script", "style", "noscript", "svg"]):
            element.decompose()
        title = soup.title.string.strip() if soup.title and soup.title.string else parsed.hostname
        text = soup.get_text("\n")
    else:
        title = parsed.hostname
        text = response.text
    extracted = _trim(text)
    if not extracted:
        raise EvidenceExtractionError("No readable public content was found at the URL.")
    return {
        "name": title,
        "url": url,
        "byte_count": len(response.content),
        "text": extracted,
    }
